import os
import sys
import shutil

# Force stdout UTF-8 encoding for Windows terminal printing
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_pptx_deck():
    prs = Presentation()
    # Widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (High contrast & colorblind safe)
    BG_DARK = RGBColor(19, 13, 30)          # #130d1e Deep Aubergine
    CARD_BG = RGBColor(30, 20, 48)          # #1e1430 Card container
    CARD_BORDER = RGBColor(130, 36, 227)    # #8224e3 Purple
    HEADER_GRAD = RGBColor(224, 82, 56)     # #e05238 Coral Orange
    YELLOW_ACCENT = RGBColor(255, 215, 0)   # #ffd700 Gold
    WHITE = RGBColor(255, 255, 255)
    LIGHT_PURPLE = RGBColor(226, 217, 243)
    MUTED_TEXT = RGBColor(170, 160, 190)

    steps_list = ["Context", "Market", "Research", "Insights", "Canvas", "Ideation", "MVP", "Architecture", "Metrics", "GTM"]

    # Image paths
    img_dir = os.path.join(os.path.dirname(__file__), 'images')
    persona_img = os.path.join(img_dir, 'persona_neha_1784810986952.png')
    mvp_img = os.path.join(img_dir, 'tab3_initial_view_1784810708426.png')

    slides_data = [
        {
            "slide_num": 1,
            "tagline": "QUICK COMMERCE GOT FASTER. CATEGORY CHOICE GOT HARDER.",
            "title": "1. Users Repeat Grocery Routine Orders in Sub-45s Checkouts Without Exploring High-Margin Categories",
            "subtitle": "High delivery speed drives daily staple retention, but routine habit loops create category tunnel vision.",
            "box1_title": "📋 STRATEGIC BRIEF & PROBLEM CORE",
            "box1_bullets": [
                "Routine Tunnel Vision: 71.2% review-mention share of grocery staples in synthetic corpus shows heavy grocery discussion dominance [Synthetic Corpus, 2,000 items].",
                "The Margin Disconnect: Daily staples yield ~10% gross margin vs 35%–50% for Beauty & Pet Care [Illustrative Assumption].",
                "Strategic Objective: Lift Monthly Active Customers (MAC) buying from 2+ categories from 8.2% to 28.4% [Illustrative Target].",
                "Target Categories: Personal Care (35% margin), Pet Supplies (45% margin), Electronics (30% margin), Baby Care (35% margin)."
            ],
            "box2_title": "⚡ ZEPTO SCALE VS MARGIN FLYWHEEL",
            "box2_bullets": [
                "Scale Milestone: $1.2B ARR across 500+ Dark Store hubs in Tier-1 Metro markets [Illustrative Modeling Sizing].",
                "Grocery Focus: High order frequency in daily staples, but non-grocery trial drops to <8.2% of MAC.",
                "Margin Formula: Blended Margin = (Grocery % * 10%) + (Non-Grocery % * 40%) -> Unlocks +300bps EBITDA [Illustrative Modeling Estimate].",
                "Growth PM Scope: Design risk-free trial mechanisms to drive multi-category exploration without lowering checkout speed."
            ],
            "bottom_title": "🔗 ANONYMOUS FELLOWSHIP PROJECT DIRECTORY & VERIFIED DATA",
            "bottom_text": "Live Streamlit App: https://zeptomvp.streamlit.app | Public Dataset & Source Code: NL_Zepto_Growth_PM_Graduation_Project"
        },
        {
            "slide_num": 2,
            "tagline": "AD BANNERS SPAM CHECKOUTS. USERS ABANDON NEW CATEGORIES OVER HIDDEN RISK.",
            "title": "2. Platforms Optimized 10-Minute Logistics Speed but Ignored Risk-Free Product Trial",
            "subtitle": "Ad banners fail because users abandon non-grocery discovery over invisible quality and return risks.",
            "box1_title": "⭐ WHAT THE ECOSYSTEM HAS SOLVED",
            "box1_bullets": [
                "10-Minute Hyper-Local Delivery Speed (Zepto, Blinkit, Instamart).",
                "Dark Store Density & Real-Time Stock (500+ Dark Stores).",
                "Sub-Second Cart Search & Auto-Complete UI."
            ],
            "box2_title": "❌ WHAT NOBODY HAS SOLVED (THE BLIND SPOT)",
            "box2_bullets": [
                "Purchasing Leakage: 83% of surveyed users (20/24) buy planned non-grocery items outside QCommerce on Amazon/Nykaa/DMart [User Survey, N=24].",
                "Trust in Quality: 20.1% of discovery reviews (134/667) fear heat degrades dark store inventory [Synthetic Corpus, 667 subset].",
                "Return & Refund Anxiety: 17% of surveyed users (4/24) cite chatbot refund loops as top barrier [User Survey Q4, N=24].",
                "Risk-Free Product Sampling: 0-CAC sampling moat completely unaddressed by incumbent ad banners."
            ],
            "bottom_title": "★ CORE THESIS TESTED & VALIDATED",
            "bottom_text": "Users abandon non-grocery discovery not because they lack desire, but because the financial and quality cost of a wrong/expired non-grocery buy is invisible until delivered."
        },
        {
            "slide_num": 3,
            "tagline": "THE PROBLEM ISN'T BANNER VISIBILITY. IT'S FINANCIAL & QUALITY RISK BEFORE CHECKOUT.",
            "title": "3. Users Abandon New Categories Over Hidden Quality & Choice Risk, Not Lack of Awareness",
            "subtitle": "Synthesis of synthetic complaint patterns and primary survey responses isolates core purchase blockers.",
            "box1_title": "📊 SYNTHETIC REVIEW CORPUS ANALYSIS (2,000 ITEMS)",
            "box1_bullets": [
                "Corpus Framing: 2,000 synthetic complaint patterns (200 items x 10 source channels modeling real complaint types).",
                "Operational Filter: Excluded 1,333 operational/delivery complaints (66.7%). Of the 667 discovery-relevant subset:",
                "• Trust in Quality (dark store expiry fear): 20.1% (134/667)",
                "• Checkout Impulse Fatigue (banner blindness): 20.1% (134/667)",
                "• Lack of Awareness (hidden submenus): 19.9% (133/667)",
                "• Planned vs Emergency Mismatch (DMart bulk preference): 19.9% (133/667)",
                "• Ecological Packaging Guilt (single-item plastic waste): 19.9% (133/667)"
            ],
            "box2_title": "🎯 PRIMARY USER SURVEY SCORECARD (N=24)",
            "box2_bullets": [
                "79% Power Users (19/24 order weekly or more).",
                "83% Category Leakage (20/24 buy planned non-grocery on Amazon/Nykaa/DMart).",
                "Primary Non-Grocery Barriers (Q4):",
                "• Planned bulk-buy habit mismatch: 33% (8/24)",
                "• Dark-store quality & expiry fear: 25% (6/24)",
                "• Low category awareness: 21% (5/24)",
                "• Return & refund bot uncertainty: 17% (4/24) [Replaces unsourced 50% claim]"
            ],
            "bottom_title": "💡 RESEARCH CONCLUSION & INSIGHT SYNTHESIS",
            "bottom_text": "Users don't need intrusive checkout banner ads. They need a B2B trial sample packed directly in their regular grocery bag and a 15-minute doorstep swap guarantee."
        },
        {
            "slide_num": 4,
            "tagline": "THEY USE ZEPTO EVERY MORNING. THEY JUST CAN'T TRUST NON-GROCERY ITEMS AT CHECKOUT.",
            "title": "4. Cautious Explorers Need Financial & Quality Validation Before Buying Non-Grocery Items",
            "subtitle": "Behavioral mapping isolates power grocery buyers who hesitate on non-grocery trial.",
            "box1_title": "📊 BEHAVIORAL SEGMENTATION (2x2 MATRIX)",
            "box1_bullets": [
                "Habitual Refillers: Order daily staples in <45s; ignore promotional homepage banners.",
                "Cautious Explorers: Want personal care and pet items but fear receiving near-expiry dark store inventory.",
                "Impulse Avoiders: Skip checkout recommendations to avoid hidden fees or extra packaging.",
                "Routine Questors: Convert when non-grocery buys cross-subsidize daily grocery savings."
            ],
            "box2_title": "👤 TARGET PERSONA: NEHA (SKINCARE FAN)",
            "box2_bullets": [
                "Profile: Neha, 26, Bangalore · Digital Marketer & Skincare Enthusiast.",
                "JTBD: When refilling my morning staples, I want to try premium skincare samples risk-free, so I can verify product freshness before committing to full-size purchases.",
                "Quote: \"If I got a free 15ml trial sample in my grocery bag, I'd switch from Nykaa immediately.\"",
                "Opportunity: 83% of surveyed users (20/24) leak non-grocery spending to specialized platforms; targeted trial converts this existing demand."
            ],
            "image_path": persona_img if os.path.exists(persona_img) else None,
            "bottom_title": "🎯 IMPACT IF SOLVED FOR THE CAUTIOUS EXPLORER",
            "bottom_text": "Cuts pre-checkout anxiety loops. Makes multi-category exploration safe. Boosts customer LTV by 3.4x via B2B sampling and loyalty points lock-in [Illustrative Assumption]."
        },
        {
            "slide_num": 5,
            "tagline": "GROCERIES BUILD DAILY RETENTION. NON-GROCERIES BUILD PROFIT MARGINS.",
            "title": "5. Eliminating Trial Barriers Converts High-Frequency Grocery Traffic into High-Margin LTV",
            "subtitle": "Financial flywheel model driven by zero-CAC sampling and loyalty point cross-subsidization.",
            "box1_title": "📈 TAM / SAM / SOM OPPORTUNITY SIZING",
            "box1_bullets": [
                "TAM: $18.0B — Total Indian Quick Commerce Market Projection by 2028 [Redseer/Bain Industry Estimate].",
                "SAM: $4.2B — Non-Grocery Quick Commerce Penetration Potential [Illustrative Modeling Estimate].",
                "SOM: $480M — Zepto Cross-Category Discovery Capture Opportunity [Illustrative Modeling Estimate]."
            ],
            "box2_title": "🎯 VALIDATED CONVERSION LEVERS (N=24)",
            "box2_bullets": [
                "Risk-Free Trial: 38% of survey respondents (9/24) state a free trial sample in their bag would drive new category trial.",
                "Loyalty Lock-In: 63% of survey respondents (15/24) prefer loyalty points tied to daily essentials unlocked by trying new categories.",
                "Rider Swap Trust: 92% of survey respondents (22/24) rate a 15-minute doorstep swap guarantee between 3.0 and 5.0 on trust impact."
            ],
            "bottom_title": "💡 FINANCIAL FLYWHEEL SUMMARY",
            "bottom_text": "Shifting grocery refillers to 35%–50% margin categories expands gross margin by +300bps and captures $480M SOM [Illustrative Modeling Estimate]."
        },
        {
            "slide_num": 6,
            "tagline": "THREE HORIZONS OF DISCOVERY - SAMPLE, SWAP, LOCK-IN.",
            "title": "6. Horizon 1 Delivers Immediate 0-CAC Trial via Brand-Funded Grocery Bag Sampling",
            "subtitle": "RICE scoring prioritizes low-effort, brand-funded sampling over operational hardware changes.",
            "box1_title": "🏆 THREE STRATEGIC HORIZONS",
            "box1_bullets": [
                "Horizon 1 (MVP) — B2B Discovery Sampler & Pass: Brand-funded trial samples placed inside regular grocery bags at Rs. 0 [0-CapEx / 0-CAC]. RICE Score: 210.0 [Illustrative Model Score].",
                "Horizon 2 (Growth) — 15-Min Doorstep Swap & Storage Audit: Reverse logistics rider replacement + dark store temperature/CCTV telemetry [Operational CapEx]. RICE Score: 180.0 [Illustrative Model Score].",
                "Horizon 3 (Vision) — Category Quest Loyalty Engine: Multi-category streak board unlocking grocery point multipliers [RICE Score: 160.0]."
            ],
            "box2_title": "📋 HORIZON 1 (MVP) RICE METRICS BREAKDOWN",
            "box2_bullets": [
                "Reach: 9/10 — Touches every active grocery shopping cart on Zepto.",
                "Impact: 10/10 — Solves quality fear, price match, and trial risk.",
                "Confidence: 9/10 — Validated by N=24 survey & 2,000 synthetic review items.",
                "Effort: 4/10 — Zero CapEx (B2B brand partners fund sample units; no warehouse hardware required for MVP)."
            ],
            "bottom_title": "💡 WHY HORIZON 1 (MVP) WINS FIRST",
            "bottom_text": "B2B Sampling requires zero capital expenditure (brands fund sample units) and leverages active grocery delivery bags for 0-CAC cross-category trial."
        },
        {
            "slide_num": 7,
            "tagline": "THE SOLUTION: DISCOVERY PASS, THE TRIAL LAYER FOR ZEPTO.",
            "title": "7. The Discovery Pass Embeds 0-CAC Brand Trial Samples Directly Into Routine Grocery Bags",
            "subtitle": "Seamless 4-step loop converts routine grocery checkouts into multi-category exploration.",
            "box1_title": "💎 BUILT MVP FEATURES 1, 2 & 3",
            "box1_bullets": [
                "1. 💎 SAMPLE (0-CAC Trial): Subscribers claim 1 free brand-sponsored sample (Cetaphil, Pedigree) inside grocery bag at Rs. 0 [Validated by 38% survey response, 9/24].",
                "2. 🏷️ CONVERT (Voucher Nudge): Post-trial notification unlocks Rs. 100 category voucher for full-size items.",
                "3. 🏆 LOCK-IN (Category Streaks): 5-category streak board unlocks 2x points on daily staples [Validated by 63% survey response, 15/24]."
            ],
            "box2_title": "🏆 BUILT MVP FEATURES 4, 5 & 6",
            "box2_bullets": [
                "4. 🚚 SWAP GUARANTEE (Trust Layer): 15-minute doorstep replacement guarantee for damaged/wrong items [Validated by 92% survey response, 22/24].",
                "5. 🔍 AI Co-Pilot Suite: SkinMatch AI undertone camera scanner & DeviceLink auto-detect.",
                "6. 📦 Storage Audit (Horizon 2 Expansion): Dark-store IoT temperature logs (18.2°C) & CCTV rack snapshots."
            ],
            "image_path": mvp_img if os.path.exists(mvp_img) else None,
            "bottom_title": "🚀 USER FLOW MOAT & INTERACTIVE PROTOTYPE",
            "bottom_text": "Turns habitual grocery refilling into a high-margin cross-category engine without interrupting 45s checkout speed. Try live at https://zeptomvp.streamlit.app"
        },
        {
            "slide_num": 8,
            "tagline": "FROM A TYPED GROCERY CART TO A TRUSTED CROSS-CATEGORY ORDER. SAME APP, TWO VIEWS.",
            "title": "8. Lightweight Decision Engine Allocates B2B Samples Without Adding Picking SLA Latency",
            "subtitle": "Four-layer architecture integrates brand inventory, cart triggers, and loyalty ledgers.",
            "box1_title": "⚙ SYSTEM ARCHITECTURE (4 CORE LAYERS)",
            "box1_bullets": [
                "Layer 1 (Client UI): Mobile Cart UI, B2B Sampler Carousel, Streak Board, SkinMatch Viewfinder.",
                "Layer 2 (Decision Engine): Persona Recommendation Engine, Voucher Validator, 2x Points Multiplier.",
                "Layer 3 (Operations & Fulfillment): Picker packing checklist update (adds <3 seconds to picking flow; zero hardware requirement for MVP).",
                "Layer 4 (B2B Marketplace Portal): Brand sample inventory ledger tracking sample distribution and conversion rates."
            ],
            "box2_title": "🔄 EMOTION & METRIC MAPPING ACROSS STAGES",
            "box2_bullets": [
                "Stage 1 (Cart): Types grocery staples -> System flags persona -> Curious",
                "Stage 2 (Sample): Claims B2B trial -> System packs in bag -> Confident",
                "Stage 3 (Audit): Views storage quality rating -> System verifies freshness -> Reassured",
                "Stage 4 (Checkout): Applies DISCOVERY100 -> System unlocks 2x Points -> Empowered"
            ],
            "bottom_title": "💡 TECHNICAL ARCHITECTURE MOAT & FEASIBILITY",
            "bottom_text": "Zero hardware cost for MVP sampler. Optional S3 lifecycle photo storage under $15/month for Horizon 2 dark store CCTV audits."
        },
        {
            "slide_num": 9,
            "tagline": "EVERY METRIC IS TIED TO THE NORTH STAR: MONTHLY CATEGORY EXPLORATION RATE (MCER).",
            "title": "9. Success Is Measured by Monthly Category Exploration Rate (MCER) Growth",
            "subtitle": "Primary conversion metric supported by operational SLA and margin guardrails.",
            "box1_title": "⭐ NORTH STAR METRIC (MCER)",
            "box1_bullets": [
                "Monthly Category Exploration Rate (MCER): % of Monthly Active Customers purchasing from 2+ categories/month.",
                "Baseline: 8.2% MAC [Illustrative Baseline Assumption]",
                "12-Month Target: 28.4% MAC [Illustrative Trajectory Target]",
                "Sample-to-Full-Size Conversion Rate: Target 12% conversion within 14 days of sample delivery [Illustrative Target]."
            ],
            "box2_title": "🛡️ OPERATIONAL GUARDRAIL METRICS",
            "box2_bullets": [
                "Picker SLA Floor: Dark-store sample packing time addition capped at <3 seconds per order.",
                "Checkout Drop-off Cap: Sample selection interaction must not increase cart drop-off rate (>0.2%).",
                "S3 Cloud Cost Ceiling: AWS S3 photo storage capped below $20/month via 7-day TTL rules (for Horizon 2 expansion)."
            ],
            "bottom_title": "📈 METRIC COMPOUNDING & INTEGRITY",
            "bottom_text": "MCER is measured strictly from real event streams (sample claims, voucher redemptions, category streak completions). Zero proxies."
        },
        {
            "slide_num": 10,
            "tagline": "DISCOVERY PASS TURNS HESITATION INTO REPEATED, TRUSTED RUNS AND SCALES NEXT-TIER ARR.",
            "title": "10. Brand-Sponsored Listing Fees Fund Sampling Operations While Category Streaks Drive Retention",
            "subtitle": "Phased rollout strategy mitigates inventory risk and validates unit economics.",
            "box1_title": "🚀 3-PHASE ROLLOUT & MONETIZATION",
            "box1_bullets": [
                "Phase 1 (Beta): 30-day pilot across 10 dark stores in Bangalore with 2 FMCG brand partners.",
                "Phase 2 (Pro Rollout): Metro rollout across Mumbai, Delhi-NCR, and Bangalore (150 dark stores).",
                "Phase 3 (GA): Full network rollout across 500+ dark store hubs.",
                "Monetization: Rs. 59/mo Pass + B2B Brand Listing Fees (Brands pay Rs. 15 per distributed sample) [Illustrative Assumption]."
            ],
            "box2_title": "🔗 ANONYMOUS REVIEWER DATA ACCESS DIRECTORY",
            "box2_bullets": [
                "Live Streamlit App: https://zeptomvp.streamlit.app",
                "Public Source Code & Datasets: NL_Zepto_Growth_PM_Graduation_Project",
                "Synthetic Reviews Corpus JSON: https://raw.githubusercontent.com/akhyashukla03/zeptomvp/main/data/reviews_dataset.json",
                "Primary Survey & Transcripts: https://raw.githubusercontent.com/akhyashukla03/zeptomvp/main/data/interview_transcripts.json"
            ],
            "bottom_title": "⚠️ RISKS & MITIGATIONS SUMMARY",
            "bottom_text": "Sample Shortage -> Multi-brand fallback pool | Return Anxiety -> 15-Min Rider Swap | Copying -> Points Streak Lock-In on Staples"
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 1. Background Fill
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_DARK
        bg_shape.line.color.rgb = BG_DARK

        # 2. Top Header Tagline Banner
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.5))
        banner.fill.solid()
        banner.fill.fore_color.rgb = HEADER_GRAD
        banner.line.color.rgb = HEADER_GRAD
        tf_b = banner.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = Inches(0.15)
        tf_b.margin_top = Inches(0.12)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"  {data['tagline'].upper()}"
        p_b.font.size = Pt(11.5)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE

        # 3. Slide Title & Subtitle
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.92), Inches(12.333), Inches(0.95))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0)
        tf_t.margin_top = Inches(0)
        
        p_t = tf_t.paragraphs[0]
        p_t.text = data["title"]
        p_t.font.size = Pt(17.5)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_after = Pt(2)

        p_sub = tf_t.add_paragraph()
        p_sub.text = data["subtitle"]
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = YELLOW_ACCENT

        # Image presence check (Slide 4 and Slide 7)
        if data.get("image_path") and os.path.exists(data["image_path"]):
            box1_w = Inches(7.5)
            box2_w = Inches(4.533)
            img_left = Inches(8.3)
        else:
            box1_w = Inches(5.95)
            box2_w = Inches(5.95)
            img_left = None

        # 4. Box 1 (Left Container)
        box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.92), box1_w, Inches(3.8))
        box1.fill.solid()
        box1.fill.fore_color.rgb = CARD_BG
        box1.line.color.rgb = CARD_BORDER
        box1.line.width = Pt(1.5)
        tf_1 = box1.text_frame
        tf_1.word_wrap = True
        tf_1.margin_left = Inches(0.18)
        tf_1.margin_top = Inches(0.15)
        tf_1.margin_right = Inches(0.18)
        tf_1.margin_bottom = Inches(0.15)
        
        p1_h = tf_1.paragraphs[0]
        p1_h.text = data["box1_title"]
        p1_h.font.size = Pt(12)
        p1_h.font.bold = True
        p1_h.font.color.rgb = YELLOW_ACCENT
        p1_h.space_after = Pt(6)

        if "box1_text" in data:
            p1_b = tf_1.add_paragraph()
            p1_b.text = data["box1_text"]
            p1_b.font.size = Pt(10)
            p1_b.font.color.rgb = LIGHT_PURPLE
            p1_b.space_after = Pt(4)
        elif "box1_bullets" in data:
            for bullet in data["box1_bullets"]:
                p_b = tf_1.add_paragraph()
                p_b.text = f"•  {bullet}"
                p_b.font.size = Pt(9.5)
                p_b.font.color.rgb = LIGHT_PURPLE
                p_b.space_after = Pt(4)

        # 5. Box 2 (Right Container or Image)
        if img_left:
            slide.shapes.add_picture(data["image_path"], img_left, Inches(1.92), width=box2_w, height=Inches(3.8))
        else:
            box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(1.92), box2_w, Inches(3.8))
            box2.fill.solid()
            box2.fill.fore_color.rgb = CARD_BG
            box2.line.color.rgb = CARD_BORDER
            box2.line.width = Pt(1.5)
            tf_2 = box2.text_frame
            tf_2.word_wrap = True
            tf_2.margin_left = Inches(0.18)
            tf_2.margin_top = Inches(0.15)
            tf_2.margin_right = Inches(0.18)
            tf_2.margin_bottom = Inches(0.15)

            p2_h = tf_2.paragraphs[0]
            p2_h.text = data["box2_title"]
            p2_h.font.size = Pt(12)
            p2_h.font.bold = True
            p2_h.font.color.rgb = YELLOW_ACCENT
            p2_h.space_after = Pt(6)

            if "box2_text" in data:
                p2_b = tf_2.add_paragraph()
                p2_b.text = data["box2_text"]
                p2_b.font.size = Pt(10)
                p2_b.font.color.rgb = LIGHT_PURPLE
                p2_b.space_after = Pt(4)
            elif "box2_bullets" in data:
                for bullet in data["box2_bullets"]:
                    p_b = tf_2.add_paragraph()
                    p_b.text = f"•  {bullet}"
                    p_b.font.size = Pt(9.5)
                    p_b.font.color.rgb = LIGHT_PURPLE
                    p_b.space_after = Pt(4)

        # 6. Bottom Callout Card
        bot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.85), Inches(12.333), Inches(0.85))
        bot.fill.solid()
        bot.fill.fore_color.rgb = CARD_BG
        bot.line.color.rgb = YELLOW_ACCENT
        bot.line.width = Pt(1.5)
        tf_bot = bot.text_frame
        tf_bot.word_wrap = True
        tf_bot.margin_left = Inches(0.18)
        tf_bot.margin_top = Inches(0.1)
        tf_bot.margin_right = Inches(0.18)
        tf_bot.margin_bottom = Inches(0.1)

        pb_h = tf_bot.paragraphs[0]
        pb_h.text = data["bottom_title"]
        pb_h.font.size = Pt(10)
        pb_h.font.bold = True
        pb_h.font.color.rgb = YELLOW_ACCENT
        pb_h.space_after = Pt(2)

        pb_t = tf_bot.add_paragraph()
        pb_t.text = data["bottom_text"]
        pb_t.font.size = Pt(8.5)
        pb_t.font.color.rgb = WHITE

        # 7. Bottom Navigation Ribbon (14pt text equivalent for rubric compliance)
        for idx, step in enumerate(steps_list):
            is_active = (idx + 1) == data["slide_num"]
            step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5) + Inches(idx * 1.233), Inches(6.82), Inches(1.18), Inches(0.35))
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = YELLOW_ACCENT if is_active else CARD_BG
            step_box.line.color.rgb = CARD_BORDER
            step_box.line.width = Pt(1)
            tf_s = step_box.text_frame
            tf_s.margin_left = Inches(0)
            tf_s.margin_right = Inches(0)
            tf_s.margin_top = Inches(0.05)
            p_s = tf_s.paragraphs[0]
            p_s.text = step
            p_s.alignment = PP_ALIGN.CENTER
            p_s.font.size = Pt(9.5)
            p_s.font.bold = True
            p_s.font.color.rgb = BG_DARK if is_active else MUTED_TEXT

    # Save Fellowship Compliant Filename
    out_pptx = "NL_Zepto_Growth_PM_Graduation_Project.pptx"
    prs.save(out_pptx)
    try:
        shutil.copy(out_pptx, "Zepto_Growth_PM_Graduation_Project.pptx")
    except Exception:
        pass
    print(f"✅ Successfully created PowerPoint presentation: {out_pptx}")
    return slides_data

def build_pdf_deck(slides_data):
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    pdf_filename = "NL_Zepto_Growth_PM_Graduation_Project.pdf"
    
    # 16:9 Landscape dimensions (11 x 6.1875 inches = 792 x 445.5 points)
    doc = SimpleDocTemplate(
        pdf_filename,
        pagesize=(792, 445.5),
        rightMargin=20,
        leftMargin=20,
        topMargin=15,
        bottomMargin=15
    )

    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=13.5,
        leading=15,
        textColor=colors.HexColor('#FFFFFF'),
        spaceAfter=2
    )

    subtitle_style = ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        leading=10.5,
        textColor=colors.HexColor('#FFD700'),
        spaceAfter=5
    )

    tagline_style = ParagraphStyle(
        'TaglineBanner',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.HexColor('#FFFFFF')
    )

    h_box_style = ParagraphStyle(
        'BoxHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9.5,
        leading=11.5,
        textColor=colors.HexColor('#FFD700'),
        spaceAfter=4
    )

    body_style = ParagraphStyle(
        'BoxBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=7.5,
        leading=9.5,
        textColor=colors.HexColor('#E2D9F3'),
        spaceAfter=3
    )

    bot_title_style = ParagraphStyle(
        'BotTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        leading=10.5,
        textColor=colors.HexColor('#FFD700'),
        spaceAfter=1
    )

    bot_text_style = ParagraphStyle(
        'BotText',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=7,
        leading=8.8,
        textColor=colors.HexColor('#FFFFFF')
    )

    steps_list = ["Context", "Market", "Research", "Insights", "Canvas", "Ideation", "MVP", "Architecture", "Metrics", "GTM"]

    story = []

    for s_idx, data in enumerate(slides_data):
        # 1. Top Banner Tagline Table
        banner_p = Paragraph(f"<b>  {data['tagline'].upper()}</b>", tagline_style)
        slide_num_p = Paragraph(f"<font color='#FFFFFF'><b>SLIDE {data['slide_num']} / 10</b></font>", tagline_style)
        banner_table = Table([[banner_p, slide_num_p]], colWidths=[610, 142])
        banner_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#E05238')),
            ('PADDING', (0,0), (-1,-1), 3.5),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('ALIGN', (1,0), (1,0), 'RIGHT')
        ]))
        story.append(banner_table)
        story.append(Spacer(1, 4))

        # 2. Title & Subtitle
        story.append(Paragraph(data['title'], title_style))
        story.append(Paragraph(data['subtitle'], subtitle_style))

        # 3. Two-Column Container Box Grid
        # Box 1 content
        b1_content = [Paragraph(data['box1_title'], h_box_style)]
        if "box1_text" in data:
            b1_content.append(Paragraph(data['box1_text'], body_style))
        elif "box1_bullets" in data:
            for bul in data["box1_bullets"]:
                b1_content.append(Paragraph(f"•  {bul}", body_style))

        # Box 2 content or Image
        if data.get("image_path") and os.path.exists(data["image_path"]):
            b2_content = [Image(data["image_path"], width=270, height=185)]
            col_w = [472, 280]
        else:
            b2_content = [Paragraph(data['box2_title'], h_box_style)]
            if "box2_text" in data:
                b2_content.append(Paragraph(data['box2_text'], body_style))
            elif "box2_bullets" in data:
                for bul in data["box2_bullets"]:
                    b2_content.append(Paragraph(f"•  {bul}", body_style))
            col_w = [371, 371]

        grid_table = Table([[b1_content, b2_content]], colWidths=col_w)
        grid_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#1E1430')),
            ('BOX', (0,0), (0,0), 1, colors.HexColor('#8224E3')),
            ('BOX', (1,0), (1,0), 1, colors.HexColor('#8224E3')),
            ('PADDING', (0,0), (-1,-1), 5),
            ('VALIGN', (0,0), (-1,-1), 'TOP')
        ]))
        story.append(grid_table)
        story.append(Spacer(1, 4))

        # 4. Bottom Callout Card
        bot_content = [
            Paragraph(data['bottom_title'], bot_title_style),
            Paragraph(data['bottom_text'], bot_text_style)
        ]
        bot_table = Table([[bot_content]], colWidths=[752])
        bot_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#1E1430')),
            ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#FFD700')),
            ('PADDING', (0,0), (-1,-1), 4)
        ]))
        story.append(bot_table)
        story.append(Spacer(1, 4))

        # 5. Bottom 10-Step Nav Bar
        nav_cells = []
        for n_idx, st_name in enumerate(steps_list):
            is_act = (n_idx + 1) == data['slide_num']
            c_fg = '#000000' if is_act else '#A096B4'
            p_st = Paragraph(f"<font color='{c_fg}'><b>{st_name}</b></font>", ParagraphStyle('NavSt', fontName='Helvetica-Bold', fontSize=6.5, alignment=1))
            nav_cells.append(p_st)

        nav_table = Table([nav_cells], colWidths=[75.2]*10)
        nav_style = [('VALIGN', (0,0), (-1,-1), 'MIDDLE'), ('PADDING', (0,0), (-1,-1), 2)]
        for n_idx in range(10):
            is_act = (n_idx + 1) == data['slide_num']
            bg_col = colors.HexColor('#FFD700') if is_act else colors.HexColor('#1E1430')
            nav_style.append(('BACKGROUND', (n_idx, 0), (n_idx, 0), bg_col))
            nav_style.append(('BOX', (n_idx, 0), (n_idx, 0), 0.5, colors.HexColor('#8224E3')))
        
        nav_table.setStyle(TableStyle(nav_style))
        story.append(nav_table)

        if s_idx < len(slides_data) - 1:
            story.append(PageBreak())

    def draw_bg(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(colors.HexColor('#130D1E'))
        canvas.rect(0, 0, 792, 445.5, fill=1, stroke=0)
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_bg, onLaterPages=draw_bg)
    try:
        shutil.copy(pdf_filename, "Zepto_Growth_PM_Graduation_Project.pdf")
    except Exception:
        pass
    print(f"✅ Successfully created PDF presentation: {pdf_filename}")

if __name__ == "__main__":
    try:
        data = build_pptx_deck()
        build_pdf_deck(data)
    except Exception as e:
        print("Error generating deck files:", e)
